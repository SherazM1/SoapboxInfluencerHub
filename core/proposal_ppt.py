from __future__ import annotations

import math
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Any


SCENARIO_KEYS = ("A", "B", "C")
CHECKMARK = "✓"
PLACEHOLDER_REPLACEMENT_SLIDE_LIMIT = 6
APPROACH_ROWS = {
    "Campaign Flight": "campaign_flight",
    "Total # of Influencers": "total_influencers",
    "Social + Story Influencers": "social_stories_count",
    "Video Influencers": "video_creators_count",
    "Total Minimum Pieces of Content": "total_minimum_pieces",
    "Click2Cart Link": "click2cart_link",
    "Paid Social": "paid_social",
    "Estimated Impressions": "estimated_impressions",
    "Estimated Engagements + Clicks": "estimated_engagements_clicks",
}


@dataclass
class ProposalResult:
    pptx_bytes: bytes
    warnings: list[str]


def number_value(value: Any) -> float:
    try:
        return float(value or 0)
    except (TypeError, ValueError):
        return 0.0


def whole_number_text(value: Any) -> str:
    numeric_value = number_value(value)
    if numeric_value <= 0:
        return ""
    return f"{int(round(numeric_value)):,}"


def format_impressions_range(value: Any) -> str:
    numeric_value = number_value(value)
    if numeric_value <= 0:
        return ""
    millions = numeric_value / 1_000_000
    floor_value = math.floor(millions)
    ceiling_value = floor_value + 1
    return f"{floor_value}MM - {ceiling_value}MM"


def format_engagements_clicks_range(value: Any) -> str:
    numeric_value = number_value(value)
    if numeric_value <= 0:
        return ""
    floor_value = math.floor(numeric_value / 50_000) * 50_000
    ceiling_value = floor_value + 50_000
    return f"{floor_value:,.0f} - {ceiling_value:,.0f}"


def format_scenario_price(value: Any) -> str:
    numeric_value = number_value(value)
    if numeric_value <= 0:
        return ""
    rounded_thousands = int(round(numeric_value / 1000))
    return f"${rounded_thousands}K"


def included_text(value: bool) -> str:
    return "true" if value else ""


def derive_total_minimum_pieces(snapshot: dict[str, Any]) -> str:
    social_stories = number_value(snapshot.get("social_stories_count"))
    video_creators = number_value(snapshot.get("video_creators_count"))
    total_pieces = social_stories * 4 + video_creators * 3
    return whole_number_text(total_pieces)


def build_scenario_payload(snapshot: dict[str, Any] | None) -> dict[str, str]:
    if not isinstance(snapshot, dict):
        return {}
    click_2_cart_cost = number_value(snapshot.get("click_2_cart_cost"))
    paid_media_spend = number_value(snapshot.get("paid_media_spend"))
    return {
        "campaign_flight": str(snapshot.get("campaign_flight") or ""),
        "total_influencers": whole_number_text(snapshot.get("total_influencers")),
        "social_stories_count": whole_number_text(snapshot.get("social_stories_count")),
        "video_creators_count": whole_number_text(snapshot.get("video_creators_count")),
        "total_minimum_pieces": derive_total_minimum_pieces(snapshot),
        "click2cart_link": included_text(click_2_cart_cost > 0),
        "paid_social": included_text(paid_media_spend > 0),
        "estimated_impressions": format_impressions_range(
            snapshot.get("organic_paid_impressions")
        ),
        "estimated_engagements_clicks": format_engagements_clicks_range(
            snapshot.get("organic_paid_engagements_clicks")
        ),
        "price": format_scenario_price(snapshot.get("program_total")),
    }


def build_proposal_payload(
    pricing_state: dict[str, Any] | None,
    scenario_snapshots: dict[str, dict[str, Any] | None] | None,
) -> dict[str, Any]:
    inputs = pricing_state.get("inputs", {}) if isinstance(pricing_state, dict) else {}
    payload = {
        "brand": inputs.get("brand") or inputs.get("client_name") or "",
        "retailer": inputs.get("retailer") or "",
        "campaign_name": inputs.get("campaign_name") or "",
        "scenarios": {},
    }
    scenario_snapshots = scenario_snapshots or {}
    for key in SCENARIO_KEYS:
        payload["scenarios"][key] = build_scenario_payload(scenario_snapshots.get(key))
    return payload


def normalize_text(value: str) -> str:
    return re.sub(r"\s+", " ", value).strip().lower()


def iter_shape_text(shape: Any) -> str:
    parts = []
    if getattr(shape, "has_text_frame", False):
        parts.append(shape.text)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                parts.append(cell.text)
    return " ".join(parts)


def iter_all_shapes(shapes: Any) -> Any:
    for shape in shapes:
        yield shape
        if hasattr(shape, "shapes"):
            yield from iter_all_shapes(shape.shapes)


def replace_paragraph_text(paragraph: Any, text: str) -> None:
    if paragraph.runs:
        first_run = paragraph.runs[0]
        for extra_run in list(paragraph.runs)[1:]:
            paragraph._p.remove(extra_run._r)
        first_run.text = text
    else:
        paragraph.add_run().text = text


def replace_text_frame_text(text_frame: Any, text: str) -> None:
    if not text_frame.paragraphs:
        text_frame.text = text
        return
    replace_paragraph_text(text_frame.paragraphs[0], text)
    for paragraph in list(text_frame.paragraphs)[1:]:
        text_frame._txBody.remove(paragraph._p)


def replace_runs_preserving_first_style(paragraph: Any, text: str) -> None:
    if paragraph.runs:
        first_run = paragraph.runs[0]
        for extra_run in list(paragraph.runs)[1:]:
            paragraph._p.remove(extra_run._r)
        first_run.text = text
    else:
        paragraph.add_run().text = text


def replacement_for(replacements: list[tuple[str, str]], target: str) -> str:
    for candidate, replacement in replacements:
        if candidate == target:
            return replacement
    return ""


def apply_placeholder_replacements(
    text: str,
    replacements: list[tuple[str, str]],
) -> str:
    updated = text
    brand = replacement_for(replacements, "Brand Name")
    retailer = replacement_for(replacements, "retailer(s)")
    if brand:
        updated = re.sub(r"\b[Bb]rand\b(?!\s+Name)(?!/product\(s\))", brand, updated)
    if retailer:
        updated = re.sub(r"\b[Rr]etailer\b(?!\(s\))", retailer, updated)

    for target, replacement in replacements:
        if replacement:
            updated = updated.replace(target, replacement)
    return updated


def replace_paragraph_placeholders(
    paragraph: Any,
    replacements: list[tuple[str, str]],
) -> bool:
    changed = False
    for run in paragraph.runs:
        original_run_text = run.text
        updated_run_text = apply_placeholder_replacements(
            original_run_text, replacements
        )
        if updated_run_text != original_run_text:
            run.text = updated_run_text
            changed = True
    if changed:
        return True

    original = "".join(run.text for run in paragraph.runs)
    if not original:
        return changed
    updated = apply_placeholder_replacements(original, replacements)
    if updated == original:
        return changed
    replace_runs_preserving_first_style(paragraph, updated)
    return True


def replace_text_frame_placeholders(
    text_frame: Any,
    replacements: list[tuple[str, str]],
) -> bool:
    changed = False
    for paragraph in text_frame.paragraphs:
        changed = replace_paragraph_placeholders(paragraph, replacements) or changed
    return changed


def replace_text_everywhere(slide: Any, replacements: list[tuple[str, str]]) -> bool:
    changed = False
    for shape in iter_all_shapes(slide.shapes):
        if getattr(shape, "has_text_frame", False):
            changed = (
                replace_text_frame_placeholders(shape.text_frame, replacements)
                or changed
            )
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    changed = (
                        replace_text_frame_placeholders(cell.text_frame, replacements)
                        or changed
                    )
    return changed


def find_slide_by_text(presentation: Any, *needles: str) -> Any | None:
    normalized_needles = [normalize_text(needle) for needle in needles]
    for slide in presentation.slides:
        slide_text = normalize_text(
            " ".join(iter_shape_text(shape) for shape in iter_all_shapes(slide.shapes))
        )
        if all(needle in slide_text for needle in normalized_needles):
            return slide
    return None


def build_template_replacements(payload: dict[str, Any]) -> list[tuple[str, str]]:
    brand = payload.get("brand") or ""
    retailer = payload.get("retailer") or ""
    campaign_name = payload.get("campaign_name") or ""
    return [
        ("Brand Name", brand),
        ("Campaign Name", campaign_name),
        (" Influencer Campaign", campaign_name),
        ("Influencer Campaign", campaign_name),
        ("brand/product(s)", brand),
        ("Brand/Product(s)", brand),
        ("retailer(s)", retailer),
        ("Retailer(s)", retailer),
        ("retailer-focused", f"{retailer}-focused" if retailer else ""),
        ("Retailer-focused", f"{retailer}-focused" if retailer else ""),
        ("retailer’s", f"{retailer}'s" if retailer else ""),
        ("Retailer’s", f"{retailer}'s" if retailer else ""),
        ("retailer shoppers", f"{retailer} shoppers" if retailer else ""),
        ("Retailer shoppers", f"{retailer} shoppers" if retailer else ""),
    ]


def populate_template_placeholders(presentation: Any, payload: dict[str, Any]) -> list[str]:
    warnings = []
    replacements = build_template_replacements(payload)
    changed = False
    for slide_index, slide in enumerate(presentation.slides, start=1):
        if slide_index > PLACEHOLDER_REPLACEMENT_SLIDE_LIMIT:
            break
        changed = replace_text_everywhere(slide, replacements) or changed
    if not changed:
        warnings.append("Proposal template placeholders were not found.")
    return warnings


def cell_key(cell_text: str) -> str:
    return normalize_text(cell_text).replace(" ", "")


def find_table_shape(slide: Any) -> Any | None:
    for shape in iter_all_shapes(slide.shapes):
        if getattr(shape, "has_table", False):
            table_text = normalize_text(iter_shape_text(shape))
            if "campaign flight" in table_text and "estimated impressions" in table_text:
                return shape
    return None


def set_cell_text(cell: Any, text: str) -> None:
    if text == "":
        return
    replace_text_frame_text(cell.text_frame, text)


def clear_cell_text(cell: Any) -> None:
    replace_text_frame_text(cell.text_frame, "")


def update_approach_slide_table(slide: Any, payload: dict[str, Any]) -> list[str]:
    table_shape = find_table_shape(slide)
    if table_shape is None:
        return ["Approach slide table was not found."]

    table = table_shape.table
    scenario_columns: dict[str, int] = {}
    for row in table.rows:
        for index, cell in enumerate(row.cells):
            text = normalize_text(cell.text)
            if text in {"a", "b", "c"}:
                scenario_columns[text.upper()] = index
    if not scenario_columns:
        return ["Approach slide A/B/C columns were not found."]

    row_map = {cell_key(label): field for label, field in APPROACH_ROWS.items()}
    price_row_index = len(table.rows) - 1
    for row_index, row in enumerate(table.rows):
        first_cell_key = cell_key(row.cells[0].text)
        field = row_map.get(first_cell_key)
        if field is None and row_index != price_row_index:
            continue
        for scenario_key, column_index in scenario_columns.items():
            scenario = payload["scenarios"].get(scenario_key, {})
            if not scenario:
                continue
            value = scenario.get("price") if row_index == price_row_index else scenario.get(field)
            if field in {"click2cart_link", "paid_social"}:
                if not value:
                    clear_cell_text(row.cells[column_index])
                continue
            if value:
                set_cell_text(row.cells[column_index], value)
    return []


def generate_powerpoint_proposal(
    template_path: str | Path,
    payload: dict[str, Any],
    output_path: str | Path,
) -> ProposalResult:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx is required to generate proposals.") from exc

    template = Path(template_path)
    output = Path(output_path)
    presentation = Presentation(str(template))
    warnings: list[str] = []

    warnings.extend(populate_template_placeholders(presentation, payload))

    approach_slide = find_slide_by_text(
        presentation,
        "Approach",
        "Campaign Flight",
        "Estimated Impressions",
    )
    if approach_slide is None:
        warnings.append("Approach slide was not found.")
    else:
        warnings.extend(update_approach_slide_table(approach_slide, payload))

    output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(output))
    return ProposalResult(pptx_bytes=output.read_bytes(), warnings=warnings)
